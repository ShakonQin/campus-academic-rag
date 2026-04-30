"""PPT文档解析器"""

from pathlib import Path
from typing import List

from pptx import Presentation
from pptx.util import Inches

from src.core.document_parser.base_parser import BaseParser, ParsedDocument
from src.utils.logger import logger


class PPTParser(BaseParser):
    """PPT文档解析器

    支持解析PPT/PPTX文档，提取幻灯片内容、备注、元数据。
    """

    SUPPORTED_EXTENSIONS = [".pptx", ".ppt"]

    def parse(self, file_path: Path) -> ParsedDocument:
        """解析PPT文档

        Args:
            file_path: PPT文件路径

        Returns:
            解析后的文档对象
        """
        if not self.validate_file(file_path):
            return self._create_error_document(file_path, "文件验证失败")

        # .ppt格式不支持（需要转换）
        if file_path.suffix.lower() == ".ppt":
            return self._create_error_document(
                file_path, "旧版.ppt格式暂不支持，请转换为.pptx格式"
            )

        try:
            # 提取基础元数据
            metadata = self._extract_metadata(file_path)

            # 解析PPT内容
            prs = Presentation(file_path)

            # 提取元数据
            core_props = prs.core_properties
            metadata["author"] = core_props.author or ""
            metadata["title"] = core_props.title or ""

            # 逐页解析
            slides = []
            full_content_parts = []
            slide_count = 0

            for slide_num, slide in enumerate(prs.slides, 1):
                try:
                    slide_content = self._extract_slide_content(slide)
                    if slide_content.strip():
                        slides.append(slide_content)
                        full_content_parts.append(
                            f"[第{slide_num}页幻灯片]\n{slide_content}"
                        )
                        slide_count += 1
                except Exception as e:
                    logger.warning(f"解析PPT第{slide_num}页失败: {e}")
                    continue

            # 构建完整内容
            full_content = "\n\n".join(full_content_parts)

            # 创建文档对象
            doc = ParsedDocument(
                doc_id=metadata["doc_id"],
                doc_name=metadata["doc_name"],
                doc_type="pptx",
                doc_path=metadata["doc_path"],
                content=full_content,
                pages=slides,
                author=metadata.get("author", ""),
                title=metadata.get("title", metadata["doc_name"]),
                created_at=metadata.get("created_at"),
                modified_at=metadata.get("modified_at"),
                page_count=slide_count,
                parse_success=True,
            )

            logger.info(f"PPT解析成功: {file_path.name}, 共{slide_count}页幻灯片")
            return doc

        except Exception as e:
            logger.error(f"PPT解析失败: {file_path.name}, 错误: {e}")
            return self._create_error_document(file_path, str(e))

    def _extract_slide_content(self, slide) -> str:
        """提取单页幻灯片内容

        Args:
            slide: 幻灯片对象

        Returns:
            幻灯片文本内容
        """
        content_parts = []

        # 提取标题
        if slide.shapes.title:
            title_text = slide.shapes.title.text.strip()
            if title_text:
                content_parts.append(f"标题: {title_text}")

        # 提取所有文本框内容
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text and text not in content_parts:
                        content_parts.append(text)

            # 提取表格内容
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_texts = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_texts.append(cell_text)
                    if row_texts:
                        content_parts.append(" | ".join(row_texts))

        # 提取备注
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                notes_text = notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    content_parts.append(f"[备注] {notes_text}")

        return "\n".join(content_parts)
